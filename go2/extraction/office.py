# Copyright (c) 2026 Mostafa Elkabir. Licensed under the BSD 2-Clause License.
"""Word and PowerPoint extraction.

Both formats carry structure worth keeping: Word headings become the
``heading`` on following blocks, and each slide stays its own block so a
citation can point at a slide number.
"""

from __future__ import annotations

import io
from typing import TYPE_CHECKING, cast

import docx
import pptx

from go2.extraction.base import Block, Extracted

if TYPE_CHECKING:
    from pptx.shapes.autoshape import Shape

_HEADING_STYLE_PREFIX = "Heading"


def extract_docx(data: bytes) -> Extracted:
    """Extract paragraphs and tables from a .docx file.

    Headings are carried onto subsequent blocks rather than emitted as blocks
    of their own, so a chunk knows which section it came from.

    Args:
        data: Raw .docx bytes.

    Returns:
        One block per paragraph, plus one block per table rendered as TSV.
    """
    document = docx.Document(io.BytesIO(data))
    blocks: list[Block] = []
    current_heading: str | None = None

    for para in document.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        if para.style is not None and str(para.style.name).startswith(_HEADING_STYLE_PREFIX):
            current_heading = text
            continue
        blocks.append(Block(text=text, heading=current_heading))

    # Word tables are usually small and prose-adjacent (unlike spreadsheets),
    # so TSV inside a block keeps them searchable without a separate artifact.
    for table in document.tables:
        rows = ["\t".join(cell.text.strip() for cell in row.cells) for row in table.rows]
        rendered = "\n".join(r for r in rows if r.strip())
        if rendered:
            blocks.append(Block(text=rendered, heading=current_heading))

    return Extracted(blocks=blocks)


def extract_pptx(data: bytes) -> Extracted:
    """Extract text from a .pptx file, one block per slide.

    Args:
        data: Raw .pptx bytes.

    Returns:
        One block per slide that contains any text, tagged with its slide number.
    """
    presentation = pptx.Presentation(io.BytesIO(data))
    blocks: list[Block] = []

    for index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            # has_text_frame is the library's own guard for this access, but it
            # is a runtime flag the checker cannot use to narrow BaseShape.
            text = cast("Shape", shape).text_frame.text.strip()
            if text:
                parts.append(text)
        if parts:
            blocks.append(Block(text="\n".join(parts), slide=index))

    return Extracted(blocks=blocks)
