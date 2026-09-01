# Copyright (c) 2026 Mostafa Elkabir. Licensed under the BSD 2-Clause License.
"""Fixture builders that produce real files, not mocks.

Every fixture returns genuine bytes in the target format, so the extractors are
exercised against the same parsers they will meet in production.
"""

from __future__ import annotations

import io

import docx
import openpyxl
import pptx
import pymupdf
import pytest
from pptx.util import Inches


@pytest.fixture
def pdf_with_text() -> bytes:
    """A two-page PDF with a real text layer on both pages."""
    doc = pymupdf.open()
    for page_number in (1, 2):
        page = doc.new_page()
        page.insert_text(
            (72, 96),
            f"Page {page_number}. The renewal fee is 4,500 USD payable to Acme Holdings.",
            fontsize=12,
        )
    data: bytes = doc.tobytes()
    doc.close()
    return data


@pytest.fixture
def pdf_scanned() -> bytes:
    """A two-page PDF with no text layer, standing in for a scan."""
    doc = pymupdf.open()
    doc.new_page()
    doc.new_page()
    data: bytes = doc.tobytes()
    doc.close()
    return data


@pytest.fixture
def pdf_mixed() -> bytes:
    """A PDF whose first page has text and whose second page does not."""
    doc = pymupdf.open()
    page = doc.new_page()
    page.insert_text((72, 96), "Contract summary and the agreed payment terms.", fontsize=12)
    doc.new_page()
    data: bytes = doc.tobytes()
    doc.close()
    return data


@pytest.fixture
def docx_bytes() -> bytes:
    """A .docx with two headed sections and a small table."""
    document = docx.Document()
    document.add_heading("Payment Terms", level=1)
    document.add_paragraph("Net 30 from invoice date.")
    document.add_heading("Termination", level=1)
    document.add_paragraph("Either party may terminate with 60 days notice.")

    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Region"
    table.cell(0, 1).text = "Quota"
    table.cell(1, 0).text = "EMEA"
    table.cell(1, 1).text = "12000"

    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


@pytest.fixture
def pptx_bytes() -> bytes:
    """A three-slide deck where the middle slide is intentionally empty."""
    presentation = pptx.Presentation()
    blank = presentation.slide_layouts[6]

    first = presentation.slides.add_slide(blank)
    first.shapes.add_textbox(
        Inches(1), Inches(1), Inches(4), Inches(1)
    ).text_frame.text = "Q3 revenue grew 18 percent"

    presentation.slides.add_slide(blank)

    third = presentation.slides.add_slide(blank)
    third.shapes.add_textbox(
        Inches(1), Inches(1), Inches(4), Inches(1)
    ).text_frame.text = "Next steps: renew the Acme contract"

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


@pytest.fixture
def xlsx_bytes() -> bytes:
    """A workbook with two sheets, one of which is empty."""
    book = openpyxl.Workbook()
    budget = book.active
    assert budget is not None
    budget.title = "Q3 Budget"
    budget.append(["Category", "Planned", "Actual"])
    budget.append(["Salaries", 120000, 118400])
    budget.append(["Cloud", 8000, 9350])

    book.create_sheet("Empty")

    buffer = io.BytesIO()
    book.save(buffer)
    return buffer.getvalue()
